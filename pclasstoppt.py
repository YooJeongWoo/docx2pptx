from slide import Slide
from pclass import Pclass
from pptx import Presentation
from pptx.util import Inches, Pt
import putImage

def makeppt(pclass):
    prs = Presentation(pclass.template_name)

    for i in range(len(pclass.Slides)):
        slide_layout = prs.slide_layouts[pclass.Slides[i].layout_num]
        temp_slide = prs.slides.add_slide(slide_layout)

        title = temp_slide.shapes.title
        title.text = pclass.Slides[i].title_text
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(55)

        text_num = len(pclass.Slides[i].additional_text)
        image_num = len(pclass.Slides[i].additional_image)
        if text_num != 0: #text가 있을 때
            for n in range(text_num):
                par = temp_slide.shapes[1].text_frame.paragraph()
                par.text = pclass.Slides[i].additional_text[n]

        if image_num != 0: #image가 있을 때
            iList = pclass.Slides[i].additional_image
            putImage.putting(len(iList), iList, temp_slide)

    prs.save('test.pptx')
    return

