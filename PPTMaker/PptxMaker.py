from PPTMaker.source import *
from Docx2PptxConverter.source.slide import Slide
from Docx2PptxConverter.source.pclass import Pclass
from pptx import Presentation

class PptxMaker:
    def make_ppt(self, pclass):
        prs = Presentation(pclass.template_name)

        for i in range(len(pclass.Slides)):
            slide_layout = prs.slide_layouts[pclass.Slides[i].layout_num]
            temp_slide = prs.slides.add_slide(slide_layout)
            temp_slide.shapes[0].text = pclass.Slides[i].title_text
            text_num = len(pclass.Slides[i].additional_text)
            image_num = len(pclass.Slides[i].additional_image)
            if text_num != 0:
                n = 0
                text = ''
                for n in range(text_num):
                    text += pclass.Slides[i].additional_text[n]

            if image_num != 0:
                n = 0
                for n in range(len(temp_slide.placeholders)):
                    if temp_slide.placeholders[n].insert_picture:
                        break
                    else:
                        n += 1
                placeholder = temp_slide.placeholders[n]
                picture = placeholder.insert_picture(pclass.Slides[i].additional_image[0])
        prs.save('test.pptx')
        return
