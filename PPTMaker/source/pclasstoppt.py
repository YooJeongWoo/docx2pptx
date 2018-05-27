from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
import PPTMaker.source.putImage as putImage
import os, sys


def makeppt(pclass):
    if getattr(sys, 'frozen', False):
        application_path = os.path.dirname(sys.executable)
    else:
        application_path = os.getcwd()
    if pclass.template_name is not None:
        prs = Presentation(str(application_path + "/template/" + pclass.template_name))
    else:
        prs = Presentation()

    for i in range(len(pclass.slides)):
        slide_layout = prs.slide_layouts[pclass.slides[i].layout_num]
        temp_slide = prs.slides.add_slide(slide_layout)

        title = temp_slide.shapes.title
        title.text = pclass.slides[i].title_text
        title.text_frame.paragraphs[0].font.bold = True
        #title.text_frame.paragraphs[0].font.size = Pt(50)
        title.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        text_num = len(pclass.slides[i].additional_text)
        image_num = len(pclass.slides[i].additional_image)
        if text_num != 0: #text가 있을 때
            for n in range(text_num):
                par = temp_slide.shapes[1].text_frame.add_paragraph()
                par.text = pclass.slides[i].additional_text[n]
                par.font.size = Pt(20)

        if image_num != 0: #image가 있을 때
            iList = pclass.slides[i].additional_image
            putImage.putting(len(iList), iList, temp_slide)

    if pclass.filename is not None:
        prs.save(str(pclass.filename) + ".pptx")
    return

